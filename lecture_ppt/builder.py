"""`Deck` 모델을 실제 .pptx 파일로 변환한다."""

from __future__ import annotations

import io

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from .models import Deck, Slide, SlideLayout

# 16:9 기본 캔버스 크기
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

ACCENT = RGBColor(0x2E, 0x5B, 0xFF)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
GRAY = RGBColor(0x55, 0x55, 0x55)


def _set_text(frame, text: str, size: int, color: RGBColor, bold: bool = False) -> None:
    frame.clear()
    p = frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _add_bullets(frame, bullets: list[str], size: int = 20) -> None:
    frame.word_wrap = True
    first = True
    for item in bullets:
        p = frame.paragraphs[0] if first else frame.add_paragraph()
        first = False
        run = p.add_run()
        run.text = item
        run.font.size = Pt(size)
        run.font.color.rgb = DARK
        p.space_after = Pt(10)


def _title_slide(prs: Presentation, deck: Deck) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 빈 레이아웃
    bar = slide.shapes.add_shape(
        1, Inches(0), Inches(3.0), SLIDE_W, Inches(0.12)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(11.5), Inches(1.2))
    _set_text(tb.text_frame, deck.title, 40, DARK, bold=True)

    if deck.subtitle:
        sb = slide.shapes.add_textbox(Inches(0.9), Inches(3.3), Inches(11.5), Inches(1.0))
        _set_text(sb.text_frame, deck.subtitle, 22, GRAY)


def _section_slide(prs: Presentation, slide_data: Slide) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DARK

    tb = slide.shapes.add_textbox(Inches(0.9), Inches(3.0), Inches(11.5), Inches(1.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = slide_data.title
    run.font.size = Pt(34)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def _content_header(slide, title: str) -> None:
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(12.0), Inches(1.0))
    _set_text(tb.text_frame, title, 28, ACCENT, bold=True)
    line = slide.shapes.add_shape(1, Inches(0.7), Inches(1.45), Inches(12.0), Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()


def _bullets_slide(prs: Presentation, slide_data: Slide) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _content_header(slide, slide_data.title)
    body = slide.shapes.add_textbox(Inches(0.9), Inches(1.9), Inches(11.5), Inches(5.0))
    _add_bullets(body.text_frame, slide_data.bullets)
    _attach_notes(slide, slide_data)


def _two_column_slide(prs: Presentation, slide_data: Slide) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _content_header(slide, slide_data.title)
    left = slide.shapes.add_textbox(Inches(0.9), Inches(1.9), Inches(5.7), Inches(5.0))
    _add_bullets(left.text_frame, slide_data.bullets)
    right = slide.shapes.add_textbox(Inches(6.9), Inches(1.9), Inches(5.7), Inches(5.0))
    _add_bullets(right.text_frame, slide_data.right_bullets)
    _attach_notes(slide, slide_data)


def _attach_notes(slide, slide_data: Slide) -> None:
    if slide_data.speaker_notes:
        slide.notes_slide.notes_text_frame.text = slide_data.speaker_notes


def build_pptx(deck: Deck) -> bytes:
    """`Deck`를 .pptx 바이트로 렌더링한다."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    for i, slide_data in enumerate(deck.slides):
        layout = slide_data.layout
        if i == 0 and layout == SlideLayout.title:
            _title_slide(prs, deck)
        elif layout == SlideLayout.title:
            # 표지 레이아웃이 본문에 오면 섹션처럼 취급
            _section_slide(prs, slide_data)
        elif layout == SlideLayout.section:
            _section_slide(prs, slide_data)
        elif layout == SlideLayout.two_column:
            _two_column_slide(prs, slide_data)
        else:
            _bullets_slide(prs, slide_data)

    # 첫 슬라이드가 표지가 아니면 앞에 표지를 하나 넣어준다.
    if not deck.slides or deck.slides[0].layout != SlideLayout.title:
        _prepend_title(prs, deck)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _prepend_title(prs: Presentation, deck: Deck) -> None:
    _title_slide(prs, deck)
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[-1])
    xml_slides.insert(0, slides[-1])
